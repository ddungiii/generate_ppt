import os
import io
import copy
from pptx import Presentation
import collections

BEFORE_DIR = "0305"
TODAY_DIR = "0312"
TEMPLATE_PPT = "template.pptx"

os.makedirs(TODAY_DIR, exist_ok=True)

"""
1. Get Template Slide
"""
template_ppt = Presentation(TEMPLATE_PPT)
if len(template_ppt.slides) > 1:
    raise Exception("❌ 템플릿 슬라이드는 1장을 초과할 수 없습니다.")
template_slide = template_ppt.slides[0]

"""
2. "{TODAY}/사진" 폴더에 있는 사진들을 기준으로, PPT를 만들 학생들 추출
"""
students = {}
for filename in os.listdir(f"{TODAY_DIR}/사진"):
    if filename.endswith(".jpg"):
        교번, 학번, 이름, _ = filename.split("_")
        key = f"{교번}_{학번}_{이름}"
        if key not in students:
            students[key] = []
        students[key].append(os.path.join(TODAY_DIR, "사진", filename))

for key in students:
    students[key].sort(
        key=lambda x: int(x.split("_")[-1].split(".")[0])  # 사진 번호로 정렬
    )

"""
3. Get image's position by Template PPT
"""
img_positions = collections.deque()
for index, shape in enumerate(template_slide.shapes):
    # 13:  Picture type, 1: Auto Shape
    # if shape.shape_type == 13 or shape.shape_type == 1:
    if shape.shape_type == 1:
        print(
            f"{shape.shape_type}, left: {shape.left.cm:.2f}cm, top: {shape.top.cm:.2f}cm"
        )
        img_positions.append(
            {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
        )

# img_positions.popleft()  # 첫 img (부산대 로고) 제외

"""
4. Generate New Slide for each Students.
"""
for student, imgs in students.items():
    try:
        input_file = os.path.join(BEFORE_DIR, "PPT", f"{student}.pptx")
        output_file = os.path.join(TODAY_DIR, "PPT", f"{student}.pptx")

        new_ppt = Presentation(input_file)
        new_slide = new_ppt.slides.add_slide(
            new_ppt.slide_layouts[6]  # 빈 슬라이드 레이아웃
        )

        # COPY
        for shape in template_slide.shapes:
            if shape.shape_type in [14, 17]:  # 텍스트 복사
                el = shape.element
                newel = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

            elif shape.shape_type == 13:  # 부산대 로고 복사
                image_stream = io.BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(
                    image_stream, shape.left, shape.top, shape.width, shape.height
                )

        for i, position in enumerate(img_positions):
            new_slide.shapes.add_picture(
                imgs[i],
                position["left"],
                position["top"],
                position["width"],
                position["height"],
            )

        new_ppt.save(output_file)
        print(f"✅ Saved: {output_file}")
    except:
        print(f"❌ failed: {output_file}")


print("🎉 All PPTs generated successfully!")
